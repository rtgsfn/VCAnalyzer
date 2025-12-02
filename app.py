import streamlit as st
import time
import os
import tempfile
from typing import Tuple, List
import pandas as pd
import docx  # pip install python-docx
from pptx import Presentation  # pip install python-pptx

# Streamlit e GUI
from streamlit_agraph import agraph, Node, Edge, Config

# LangChain - Documenti e RAG
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings

import io
from fpdf import FPDF # Per il PDF


try:
    from analyzer import AgenticKRAG
    from extractor import KnowledgeGraph, DocumentClaims, Claim, CompetitorAnalysis
    from graph import GraphTool
    from scraper import scrape_article_text
    from vc_metrics import VCMetricsProfile, PharmaMetricsProfile, REMetricsProfile, LegalMetricsProfile # Import per la tipizzazione corretta
except ImportError as e:
    st.error(f"❌ Errore Critico: Impossibile importare i moduli. Dettagli: {e}")
    st.stop()

# ============================================================================
# CONFIGURAZIONE
# ============================================================================

MODEL_OPTIONS = {
    "Groq": [
        "llama-3.3-70b-versatile",
        "llama-3.1-8b-instant",
        "openai/gpt-oss-20b",
        "openai/gpt-oss-120b",
        "qwen/qwen3-32b"
    ],
    "Perplexity": [
        "llama-3-8b-instruct",
        "llama-3-70b-instruct",
        "mixtral-8x7b-instruct"
    ],
    "Google": [
        "gemini-2.5-pro",
        "gemini-2.5-flash",
        "gemini-2.0-flash"
    ],
    "Ollama": [
        "llama3:8b-instruct",
        "gemma:7b-instruct",
        "mixtral:8x7b-instruct",
        "kimi-k2:1t-cloud"
    ]
}


# ============================================================================
# FUNZIONI DI UTILITÀ (ESTRAZIONE TESTO)
# ============================================================================
def sanitize_url(url: str) -> str:
    """Assicura che l'URL abbia il protocollo per evitare reload della pagina."""
    if not url:
        return ""
    if not url.startswith(("http://", "https://")):
        return f"https://{url}"
    return url

def inject_custom_css():
    st.markdown("""
    <style>
        /* --- UNCERTAINTY VISUALIZATION --- */
        .metric-card {
            background-color: #f9f9f9;
            border-radius: 8px;
            padding: 15px;
            margin-bottom: 10px;
            box-shadow: 0 2px 5px rgba(0,0,0,0.05);
            transition: transform 0.2s;
        }
        .conf-high { border-left: 5px solid #2ECC71; }
        .conf-low { 
            border-left: 5px solid #F39C12;
            border: 1px dashed #F39C12;
            background-color: #fffcf5;
        }
        .metric-value { font-size: 24px; font-weight: bold; color: #2C3E50; }
        .metric-label { font-size: 14px; color: #7F8C8D; text-transform: uppercase; }
        .conf-badge {
            font-size: 10px; padding: 2px 6px; border-radius: 4px; float: right; font-weight: bold;
        }
        .badge-high { background-color: #d4edda; color: #155724; }
        .badge-low { background-color: #fff3cd; color: #856404; }

        /* --- CONTEXT HIGHLIGHTING --- */
        .source-chunk {
            font-size: 12px; background-color: #eef; border-left: 3px solid #4a90e2;
            padding: 10px; margin-bottom: 10px; font-family: monospace;
        }
        .source-header { font-weight: bold; color: #4a90e2; margin-bottom: 4px; }

        /* Pulsanti (cerchi) */
        div.vis-network div.vis-navigation div.vis-button {
            background-color: #FFFFFF !important; /* Sfondo bianco per contrasto */
            border-radius: 50% !important;
            border: 2px solid #000000 !important; /* Bordo nero spesso */
            box-shadow: 0 2px 5px rgba(0,0,0,0.2) !important;
        }
        
        /* Hover sui pulsanti */
        div.vis-network div.vis-navigation div.vis-button:hover {
            background-color: #E0E0E0 !important; /* Grigio chiaro al passaggio */
            box-shadow: 0 4px 8px rgba(0,0,0,0.3) !important;
        }

        /* Icone interne (le frecce dentro i cerchi) */
        div.vis-network div.vis-navigation div.vis-button:before {
            color: #000000 !important; /* Frecce NERE */
        }

        /* --- LABEL DEGLI ARCHI (TESTO NERO) --- */
        div.vis-network div.vis-label {
            color: #000000 !important; /* Forza il testo delle label a nero */
            font-weight: bold !important; /* Opzionale: rende il testo più leggibile */
        }
    </style>
    """, unsafe_allow_html=True)

def render_competitor_tab(comp_analysis: CompetitorAnalysis):
    """Renderizza il tab dei competitor."""
    if not comp_analysis or not comp_analysis.competitors:
        st.info("ℹ️ Nessun dato sui competitor disponibile per questa analisi.")
        return

    st.markdown("### ⚔️ Competitive Landscape")
    st.info(f"**📍 Market Position:** {comp_analysis.market_position}")

    for comp in comp_analysis.competitors:
        with st.expander(f"🏢 {comp.name}", expanded=True):
            cols = st.columns([3, 1])
            with cols[0]:
                st.markdown(f"**Cosa fanno:** {comp.description}")
                if comp.differentiation:
                    st.markdown(f"👉 **Confronto:** *{comp.differentiation}*")
            with cols[1]:
                if comp.website:
                    # FIX: Sanitizziamo l'URL
                    clean_url = sanitize_url(comp.website)
                    st.link_button("🌐 Visita Sito", clean_url)
                else:
                    st.caption("No URL")

def render_confidence_metric(label: str, value: str, status: str = "VERIFIED", explanation: str = ""):
    """
    Renderizza una metrica con visualizzazione dell'incertezza.
    Versione FIXATA per evitare il bug del </div> visibile.
    """
    # Logica di confidenza
    # Mappiamo anche i valori italiani o nulli per sicurezza
    safe_status = str(status).upper() if status else "MISSING"

    if safe_status in ["VERIFIED", "VERIFICATA", "EXCELLENT", "GOOD"]:
        css_class = "conf-high"
        badge_class = "badge-high"
        badge_text = "HIGH CONFIDENCE"
        icon = "✅"
    else:
        # UNVERIFIED, MISSING, POOR, CONFLICTING -> Low Confidence
        css_class = "conf-low"
        badge_class = "badge-low"
        badge_text = "NEEDS VERIFICATION"
        icon = "⚠️"

    # HTML compatto senza indentazione per evitare errori di rendering
    html_code = f"""
<div class="metric-card {css_class}">
    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 5px;">
        <span class="metric-label">{label}</span>
        <span class="conf-badge {badge_class}">{icon} {badge_text}</span>
    </div>
    <div class="metric-value">{value}</div>
    <div style="font-size: 11px; color: #666; margin-top: 5px; min-height: 15px;">{explanation}</div>
</div>
"""
    st.markdown(html_code, unsafe_allow_html=True)

def render_source_inspector(source_docs: list):
    """Renderizza i chunk RAG originali nella sidebar (Ground Truth)."""
    with st.sidebar:
        st.markdown("---")
        st.markdown("### 🔍 Context Inspector")
        st.caption("Verifica qui i frammenti originali usati dall'AI.")

        if not source_docs:
            st.info("Nessun contesto specifico recuperato.")
            return

        for i, doc in enumerate(source_docs):
            # Estrae il nome file dai metadati o usa un default
            source_name = doc.metadata.get("source", "Documento sconosciuto")
            page_num = doc.metadata.get("page", "?")

            with st.expander(f"📄 Fonte {i + 1}: {os.path.basename(source_name)} (p.{page_num})"):
                # Evidenziazione visiva del chunk
                st.markdown(f"""
                <div class="source-chunk">
                    <div class="source-header">{source_name}</div>
                    {doc.page_content}
                </div>
                """, unsafe_allow_html=True)

def extract_text_from_file(uploaded_file) -> str:
    """Estrae testo da PDF, DOCX, PPTX, TXT. Usato per i Requisiti."""
    text = ""
    suffix = os.path.splitext(uploaded_file.name)[1].lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_file_path = tmp_file.name

    try:
        if suffix == '.pdf':
            loader = PyPDFLoader(tmp_file_path)
            docs = loader.load()
            text = "\n".join([d.page_content for d in docs])
        elif suffix == '.txt':
            loader = TextLoader(tmp_file_path, encoding='utf-8')
            docs = loader.load()
            text = docs[0].page_content
        elif suffix in ['.docx', '.doc']:
            doc = docx.Document(tmp_file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif suffix in ['.pptx', '.ppt']:
            prs = Presentation(tmp_file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            text = "\n".join(text_runs)
    except Exception as e:
        st.error(f"❌ Errore lettura file {uploaded_file.name}: {e}")
    finally:
        if os.path.exists(tmp_file_path):
            os.remove(tmp_file_path)

    return text

def render_requirements_sidebar(report_text: str, metadata: dict):
    """
    Renderizza una sidebar specifica per la Gap Analysis (Matching Requisiti).
    Calcola statistiche basandosi sulle emoji presenti nel report.
    """
    with st.sidebar:
        st.markdown("---")
        st.markdown("### 🛠️ Gap Analysis Dashboard")

        # Parsing euristico (conta le occorrenze delle icone nel testo markdown)
        # Nota: Questo funziona perché l'LLM è istruito a usare queste emoji specifiche
        n_satisfied = report_text.count("✅")
        n_partial = report_text.count("⚠️")
        n_missing = report_text.count("❌")

        total_reqs = n_satisfied + n_partial + n_missing

        if total_reqs > 0:
            # Calcolo score di copertura (OK = 100%, Partial = 50%, Missing = 0%)
            score_pct = (n_satisfied + (0.5 * n_partial)) / total_reqs

            st.metric("Copertura Stimata", f"{score_pct:.0%}")
            st.progress(score_pct)

            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("✅ OK", n_satisfied)
            with col2:
                st.metric("⚠️ Parz.", n_partial)
            with col3:
                st.metric("❌ Gap", n_missing)

        else:
            st.warning("Nessun requisito tracciato nel report.")

        st.markdown("---")
        st.markdown("### ⚡ Performance")
        st.metric("⏱️ Tempo Analisi", f"{metadata.get('analysis_time_seconds', 0):.1f}s")

        st.info("""
        **Legenda:**
        ✅ Requisito Soddisfatto
        ⚠️ Copertura Parziale
        ❌ Gap Rilevato / Non Trovato
        """)

# ============================================================================
# FUNZIONI DI CACHE
# ============================================================================

@st.cache_resource(show_spinner="🔧 Inizializzazione agente...")
def get_agent_instance(provider: str, model_name: str) -> AgenticKRAG:
    """Inizializza e mette in cache l'istanza dell'Agente KRAG."""
    try:
        agent = AgenticKRAG(provider=provider, model_name=model_name)
        return agent
    except Exception as e:
        st.error(f"❌ Errore critico durante l'inizializzazione: {e}")
        return None


@st.cache_resource(show_spinner="📦 Caricamento modello embedding...")
def load_embedding_resources():
    """Carica e mette in cache le risorse per il RAG."""
    embeddings = HuggingFaceEmbeddings(
        model_name="all-MiniLM-L6-v2",
        model_kwargs={'device': 'cpu'}
    )
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    return embeddings, splitter


@st.cache_resource(show_spinner="📄 Elaborazione documenti (PDF, TXT, DOCX, PPTX)...")
def process_documents(files_or_text, _embeddings, _splitter) -> Tuple[object, str]:
    """Processa file eterogenei e restituisce (retriever, full_text)."""
    all_documents = []
    full_text_list = []

    if isinstance(files_or_text, str):  # Testo incollato
        full_text_list.append(files_or_text)
        all_documents = [Document(page_content=files_or_text, metadata={"source": "pasted_text"})]

    elif isinstance(files_or_text, list):  # Lista di file caricati
        for uploaded_file in files_or_text:
            # Usiamo la funzione helper per estrarre il testo grezzo
            extracted_text = extract_text_from_file(uploaded_file)

            if extracted_text:
                # --- MODIFICA IMPORTANTE: INSERIAMO IL NOME DEL FILE NEL TESTO ---
                # Questo permette all'LLM di citare la fonte nel report
                text_with_header = f"\n\n--- FILE: {uploaded_file.name} ---\n{extracted_text}"

                full_text_list.append(text_with_header)

                # Creiamo il Document object per LangChain
                all_documents.append(Document(
                    page_content=extracted_text,
                    metadata={"source": uploaded_file.name}
                ))

    if not all_documents:
        st.error("❌ Nessun documento valido processato")
        return None, None

    # Uniamo tutto il testo
    full_text = "\n".join(full_text_list)

    chunks = _splitter.split_documents(all_documents)
    vectordb = Chroma.from_documents(documents=chunks, embedding=_embeddings)

    return vectordb.as_retriever(search_kwargs={"k": 5}), full_text


def make_status_callback(container):
    """Wrapper per adattare (message, level) alle chiamate Streamlit."""

    def _callback(message: str, level: str = "info"):
        if level == "error":
            container.error(message)
        elif level == "warning":
            container.warning(message)
        elif level == "success":
            container.success(message)
        else:
            container.info(message)

    return _callback


def render_fact_checking_table(claims_data: List[dict]):
    """Renderizza tabella fact-checking con LINK CLICCABILI."""
    if not claims_data:
        st.info("ℹ️ Nessuna affermazione fattuale estratta")
        return

    # Creazione DataFrame
    df = pd.DataFrame(claims_data)

    # Mappatura icone status
    status_icons = {
        "VERIFICATA": "✅",
        "FALSA": "❌",
        "PARZIALMENTE VERIFICATA": "⚠️",
        "NON VERIFICABILE": "❓"
    }

    # 1. Colonna Status abbellita
    df["Status"] = df["status"].apply(lambda x: f"{status_icons.get(x, '?')} {x}")

    # 2. Gestione URL: Sanitizziamo e gestiamo i vuoti
    # Assicuriamoci che la chiave 'source_url' esista (per retrocompatibilità)
    if "source_url" not in df.columns:
        df["source_url"] = None

    df["url_clean"] = df["source_url"].apply(sanitize_url)

    # 3. Preparazione DataFrame per la visualizzazione
    # Selezioniamo le colonne da mostrare. 'url_clean' è quella che diventerà il link.
    df_display = df[["soggetto", "affermazione", "Status", "prove", "url_clean"]].copy()
    df_display.columns = ["Soggetto", "Affermazione", "Verdetto", "Prove", "Fonte"]

    # Statistiche (Contatori sopra la tabella)
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("✅ Verificate", len([c for c in claims_data if c["status"] == "VERIFICATA"]))
    with col2:
        st.metric("❌ False", len([c for c in claims_data if c["status"] == "FALSA"]))
    with col3:
        st.metric("⚠️ Parziali", len([c for c in claims_data if c["status"] == "PARZIALMENTE VERIFICATA"]))
    with col4:
        st.metric("❓ Non Verif.", len([c for c in claims_data if c["status"] == "NON VERIFICABILE"]))

    # 4. Render Tabella con Configurazione Link
    st.dataframe(
        df_display,
        use_container_width=True,
        hide_index=True,
        column_config={
            "Soggetto": st.column_config.TextColumn("Soggetto", width="small"),
            "Affermazione": st.column_config.TextColumn("Claim", width="medium"),
            "Verdetto": st.column_config.TextColumn("Esito", width="small"),
            "Prove": st.column_config.TextColumn("Evidenza trovata", width="large"),

            # CONFIGURAZIONE CHIAVE PER I LINK
            "Fonte": st.column_config.LinkColumn(
                "Fonte Web",
                display_text="🔗 Apri Fonte",  # Questo è il testo che appare cliccabile
                help="Clicca per visitare la pagina originale usata per la verifica",
                width="small"
            )
        }
    )


def render_knowledge_graph(nodes_data, edges_data, entity_name):
    """Renderizza grafo con styling migliorato."""
    if not nodes_data:
        st.info(f"ℹ️ Nessun dato nel Knowledge Graph per '{entity_name}'")
        return

    # Usiamo la palette centralizzata da config.py se disponibile, altrimenti fallback
    try:
        from config import GRAPH_COLOR_PALETTE, GRAPH_CONFIG
        color_map = GRAPH_COLOR_PALETTE
        # Creiamo un oggetto Config di agraph usando il dizionario
        # Nota: Streamlit-agraph vuole un oggetto Config, non un dict diretto
        config = Config(
            width=GRAPH_CONFIG["width"],
            height=GRAPH_CONFIG["height"],
            directed=GRAPH_CONFIG["directed"],
            physics=GRAPH_CONFIG["physics"],
            hierarchical=False,
            # Aggiungiamo opzioni extra per la UI
            interaction={"navigationButtons": True, "zoomView": True}
        )
    except ImportError:
        # Fallback se config non è importabile
        color_map = {"Startup": "#FF6B6B", "Persona": "#4ECDC4", "Investitore": "#45B7D1", "Entita": "#95E1D3"}
        config = Config(width=900, height=600, directed=True, physics=True, hierarchical=False)
    nodes = []
    seen_nodes = set()
    for n in nodes_data:
        if n["id"] in seen_nodes:
            continue
        seen_nodes.add(n["id"])
        # Determina il tipo e il colore
        node_type = n["label"].split("\n")[0] if "\n" in n["label"] else "Entita"

        # Logica speciale per il nodo centrale (Focus)
        if n["id"] == entity_name:
            color = color_map.get("focus_entity", "#FFD93D")
            size = 40  # Nodo centrale ben visibile
            font_size = 20
        else:
            color = color_map.get(node_type, "#95E1D3")
            size = 20  # Nodi periferici standard
            font_size = 14

        # Aggiungiamo il nodo con le proprietà fisiche
        nodes.append(Node(
            id=n["id"],
            label=n["label"],
            color=color,
            size=size,
            font={"size": font_size, "color": "#2C3E50"},
            # Opzionale: Shape diversa per tipo
            shape="dot"
        ))

    # Costruzione Archi
    edges = []
    for e in edges_data:
        # Colora l'arco in base al tipo di relazione (se presente nella mappa)
        # es. HA_FONDATO -> Verde, HA_INVESTITO -> Blu
        rel_type = e["label"].split(" ")[0]  # Prende la prima parola della label
        edge_color = color_map.get(rel_type, color_map.get("default_edge", "#A9D0F5"))  # Grigio default

        edges.append(Edge(
            source=e["source"],
            target=e["target"],
            label=e["label"],
            color=edge_color,
            width=2,
            # Aggiungiamo frecce per la direzione
            arrows="to"
        ))

    st.markdown(f"### 🕸️ Knowledge Graph per **{entity_name}**")
    st.caption("Usa la rotellina per zoomare e trascina per navigare.")

    # Render finale
    agraph(nodes=nodes, edges=edges, config=config)


def render_metadata_sidebar(metadata: dict):
    """Renderizza sidebar con metadata analisi (Versione Aggiornata Tesi)."""
    with st.sidebar:
        st.markdown("---")
        st.markdown("### 🌟 Dashboard Affidabilità")

        if "error" in metadata:
            st.error(f"❌ {metadata['error']}")
            return

        # Recuperiamo il nuovo score composito
        rel_score = metadata.get('confidence_fact_check_score', 0)

        # Determina colore e label in base al punteggio
        if rel_score >= 0.7:
            score_color = "green"
            score_label = "Alta"
        elif rel_score >= 0.4:
            score_color = "orange"
            score_label = "Media"
        else:
            score_color = "red"
            score_label = "Bassa"

        st.metric("🛡️ Indice Affidabilità (Composite)", f"{rel_score:.0%}", score_label)
        st.progress(rel_score)

        st.caption(f"""
        **Composizione Score:**
        - ✅ Veridicità Fattuale (50%)
        - 🕸️ Corroborazione Grafo (30%)
        - 🎯 Confidenza Entità (20%)
        """)

        st.markdown("---")

        # Mostriamo i dettagli grezzi
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Fatti Verificati", metadata.get('claims_verified', 0))
        with col2:
            st.metric("Nodi Grafo", metadata.get('graph_nodes', 0))

        st.metric("⏱️ Tempo Analisi", f"{metadata.get('analysis_time_seconds', 0):.1f}s")


# ============================================================================
# INTERFACCIA PRINCIPALE
# ============================================================================

st.set_page_config(page_title="VC KRAG Analyzer", page_icon="🎯", layout="wide")

st.title("SPECTRE - Agentic Intelligence")
st.markdown("**Knowledge-Retrieval Augmented Generation** per Due Diligence Investigativa")
st.markdown("---")

# Sidebar Config
with st.sidebar:
    st.header("⚙️ Configurazione")
    sector = st.selectbox(
        "🏢 Settore di Analisi",
        ["Venture Capital", "Real Estate", "Legal / M&A", "Pharma & Biotech"],
        index=0,
        help="Cambia la prospettiva dell'agente (es. Analista Finanziario vs Avvocato)"
    )

    st.divider()
    provider = st.selectbox("🤖 Provider LLM:", list(MODEL_OPTIONS.keys()))
    available_models = MODEL_OPTIONS.get(provider, [])
    model_name = st.selectbox("📦 Modello:", available_models)

try:
    embeddings, splitter = load_embedding_resources()
except Exception as e:
    st.error(f"❌ Errore caricamento risorse: {e}")
    st.stop()

# 1. Input Documenti Principali
st.markdown("## 📂 Step 1: Carica Documenti da Analizzare")
st.caption("Supporta PDF, DOCX, PPTX, TXT")
tab1, tab2 = st.tabs(["📁 Upload Files", "📋 Paste Text"])

with tab1:
    uploaded_files = st.file_uploader(
        "Carica documenti (Pitch deck, Memo, etc.)",
        type=["txt", "pdf", "docx", "doc", "pptx", "ppt"],
        accept_multiple_files=True
    )
with tab2:
    pasted_text = st.text_area("Incolla testo del documento", height=150)

st.markdown("---")

# 2. Configurazione Analisi
st.markdown("## 🎯 Step 2: Configurazione Analisi")

col_conf1, col_conf2 = st.columns([1, 2])

# Variabili di default (per evitare errori se non definite)
entita_utente = ""
requirements_text = None
is_deep_search = False
use_auto_detect = False

with col_conf1:
    analysis_mode = st.radio(
        "Tipo di Analisi:",
        ["Standard Due Diligence", "Matching Requisiti Tecnici"],
        help="Standard: Analisi rischi/opportunità. Requisiti: Verifica puntuale di una lista di necessità."
    )

# Logica Condizionale per la UI
if analysis_mode == "Matching Requisiti Tecnici":
    # --- MODALITÀ MATCHING (Requisiti visibili, Entità/Web nascosti) ---
    st.info("📋 **Modalità Requisiti**: Carica la lista delle necessità tecniche/funzionali.")
    req_tab1, req_tab2 = st.tabs(["📂 Carica File Requisiti", "✍️ Incolla Lista"])

    with req_tab1:
        req_file = st.file_uploader(
            "Carica file requisiti (.pdf, .docx, .pptx, .txt)",
            type=["txt", "pdf", "docx", "doc", "pptx", "ppt"],
            key="req_uploader"
        )
        if req_file:
            requirements_text = extract_text_from_file(req_file)
            if requirements_text:
                st.success(f"✅ Requisiti caricati ({len(requirements_text)} caratteri)")

    with req_tab2:
        req_input = st.text_area("Incolla qui la lista tecnica:", height=150,
                                 placeholder="- Supporto SSO SAML\n- Hosting On-Premise...")
        if req_input:
            requirements_text = req_input

    # Nella colonna 2 non mostriamo nulla in questa modalità
    with col_conf2:
        st.write("")

else:
    # --- MODALITÀ STANDARD (Entità/Web visibili, Requisiti nascosti) ---
    with col_conf2:
        entita_utente = st.text_input("Entità Target (Opzionale se rilevabile):", placeholder="Es. 'Figure AI'")
        use_auto_detect = st.checkbox("🤖 Rileva automaticamente l'entità", value=True)
        is_deep_search = st.toggle("🚀 Deep Search (Web)", value=False)

st.markdown("---")

# ============================================================================
# ESECUZIONE ANALISI (UNICO PULSANTE)
# ============================================================================
# ============================================================================
# ESECUZIONE ANALISI (LOGICA PERSISTENTE)
# ============================================================================

# 1. LOGICA DI ESECUZIONE (Click Bottone)
if st.button("🚀 Avvia Analisi", type="primary", width='stretch'):

    # Validazione Input
    input_data = uploaded_files if uploaded_files else pasted_text
    if not input_data:
        st.error("❌ Manca il documento da analizzare (Step 1)")
        st.stop()

    if analysis_mode == "Matching Requisiti Tecnici" and not requirements_text:
        st.error("❌ Manca la lista dei requisiti per il matching.")
        st.stop()

    if analysis_mode == "Standard Due Diligence" and not entita_utente and not use_auto_detect:
        st.error("❌ Specifica un'entità da analizzare o abilita il rilevamento automatico.")
        st.stop()

    # Reset stato precedente
    if "analysis_results" in st.session_state:
        del st.session_state["analysis_results"]

    status_container = st.empty()
    progress_bar = st.progress(0)

    try:
        # Inizializzazione Agente
        status_container.info("🔧 Inizializzazione agente...")
        agent = get_agent_instance(provider, model_name)
        if not agent: st.stop()
        agent.status_callback = make_status_callback(status_container)

        # Processamento Documenti
        progress_bar.progress(10)
        status_container.info("📄 Processamento documenti...")
        doc_retriever, document_text = process_documents(input_data, embeddings, splitter)
        if not doc_retriever or not document_text: st.stop()

        # Deduzione Entità
        entity_to_analyze = entita_utente.strip()
        should_deduce = False

        if analysis_mode == "Matching Requisiti Tecnici":
            entity_to_analyze = "Gap Analysis"
        elif use_auto_detect and not entity_to_analyze:
            should_deduce = True

        if should_deduce:
            status_container.info("🔍 Deduzione entità in corso...")
            deduction = agent.deduce_entity_from_document(document_text)
            if deduction["entity_found"]:
                entity_to_analyze = deduction["entity_name"]
                status_container.success(f"✅ Identificato: {entity_to_analyze}")
            else:
                st.error("Impossibile identificare l'entità automaticamente.")
                st.stop()

        # Avvio Analisi
        start_time = time.time()
        progress_bar.progress(20)
        actual_deep_search = False if analysis_mode == "Matching Requisiti Tecnici" else is_deep_search

        result = agent.run_full_analysis_streaming(
            entity_to_analyze,
            document_text,
            doc_retriever,
            is_deep_search=actual_deep_search,
            requirements_text=requirements_text,
            sector=sector
        )

        if "error" in result:
            st.error(f"❌ Errore: {result['error']}")
            st.stop()

        progress_bar.progress(100)
        elapsed = time.time() - start_time
        status_container.empty()

        # --- SALVATAGGIO IN SESSION STATE (SOLUZIONE AL CRASH) ---
        st.session_state["analysis_results"] = result
        st.session_state["analysis_metadata"] = {
            "entity": entity_to_analyze,
            "elapsed": elapsed,
            "mode": analysis_mode,
            "sector": sector,
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S")
        }

        # Ricarica la pagina per attivare la modalità visualizzazione
        st.rerun()

    except Exception as e:
        st.error(f"❌ Errore durante l'esecuzione: {e}")
        import traceback

        with st.expander("Dettagli tecnico"):
            st.code(traceback.format_exc())

# 2. LOGICA DI VISUALIZZAZIONE (Persistente)
# Questo blocco viene eseguito ad ogni ricaricamento (zoom, pan, click) se ci sono dati.
if "analysis_results" in st.session_state:
    result = st.session_state["analysis_results"]
    meta = st.session_state["analysis_metadata"]

    entity_to_analyze = meta["entity"]
    analysis_mode = meta["mode"]
    sector = meta["sector"]
    elapsed = meta["elapsed"]

    # Placeholders
    metrics_placeholder = st.empty()
    results_placeholder = st.empty()
    graph_placeholder = st.empty()

    # A) MODALITÀ MATCHING REQUISITI
    if analysis_mode == "Matching Requisiti Tecnici":
        with results_placeholder.container():
            st.markdown(f"# 📋 Report Matching Requisiti")
            st.caption(f"Analisi completata in {elapsed:.1f} secondi su documenti interni.")
            st.markdown("---")
            st.markdown(result.get("executive_summary", ""))

    # B) MODALITÀ STANDARD (VC DUE DILIGENCE)
    else:
        with metrics_placeholder.container():
            st.markdown(f"# 🚀 {sector} Analysis Report")
            #st.caption(f"Analisi completata in {elapsed:.1f} secondi")
            st.markdown("---")

            metrics = result.get("metrics")
            if metrics:
                st.markdown(f"## 📊 Metriche Chiave Recuperate")

                # VC PROFILE
                if isinstance(metrics, VCMetricsProfile):
                    tab_names = ["💰 SaaS", "📈 Traction", "🌍 Market", "👥 Team", "💵 Fundraising"]
                    tab_saas, tab_traction, tab_market, tab_team, tab_fundraising = st.tabs(tab_names)

                    with tab_saas:
                        m = metrics.saas_metrics
                        if m:
                            statuses = getattr(m, "metrics_status", {}) or {}
                            c1, c2, c3 = st.columns(3)
                            with c1:
                                render_confidence_metric("ARR", f"${m.arr}M" if m.arr else "-",
                                                         statuses.get("arr", "UNVERIFIED"))
                                render_confidence_metric("Growth YoY",
                                                         f"{m.revenue_growth_rate}%" if m.revenue_growth_rate else "-",
                                                         statuses.get("revenue_growth_rate", "UNVERIFIED"))
                            with c2:
                                render_confidence_metric("LTV/CAC", f"{m.ltv_cac_ratio}x" if m.ltv_cac_ratio else "-",
                                                         statuses.get("ltv_cac_ratio", "UNVERIFIED"))
                                render_confidence_metric("Net Retention",
                                                         f"{m.net_retention_rate}%" if m.net_retention_rate else "-",
                                                         statuses.get("net_retention_rate", "UNVERIFIED"))
                            with c3:
                                render_confidence_metric("Runway", f"{m.runway_months} mo" if m.runway_months else "-",
                                                         statuses.get("runway_months", "UNVERIFIED"))
                        else:
                            st.info("Nessuna metrica SaaS.")

                    with tab_traction:
                        m = metrics.traction_metrics
                        if m:
                            statuses = getattr(m, "metrics_status", {}) or {}
                            c1, c2 = st.columns(2)
                            with c1:
                                render_confidence_metric("Utenti Totali",
                                                         f"{m.total_users:,}" if m.total_users else "-",
                                                         statuses.get("total_users", "UNVERIFIED"))
                            with c2:
                                render_confidence_metric("NPS Score", str(m.nps_score) if m.nps_score else "-",
                                                         statuses.get("nps_score", "UNVERIFIED"))
                        else:
                            st.info("Nessuna metrica Traction.")

                    with tab_market:
                        m = metrics.market_metrics
                        if m:
                            statuses = getattr(m, "metrics_status", {}) or {}
                            c1, c2, c3 = st.columns(3)
                            with c1:
                                render_confidence_metric("TAM", f"${m.tam}B" if m.tam else "-",
                                                         statuses.get("tam", "UNVERIFIED"))
                            with c2:
                                render_confidence_metric("SAM", f"${m.sam}B" if m.sam else "-",
                                                         statuses.get("sam", "UNVERIFIED"))
                            with c3:
                                render_confidence_metric("SOM", f"${m.som}M" if m.som else "-",
                                                         statuses.get("som", "UNVERIFIED"))
                        else:
                            st.info("Nessuna metrica Market.")

                    with tab_fundraising:
                        if metrics.fundraising_metrics and metrics.fundraising_metrics.rounds:
                            for round_data in metrics.fundraising_metrics.rounds:
                                st.markdown(f"**{round_data.stage.value}** - ${round_data.amount}M")
                                if round_data.lead_investor: st.caption(f"Lead: {round_data.lead_investor}")
                                st.markdown("---")
                        else:
                            st.info("Nessun fundraising.")

                    if metrics.team_metrics and metrics.team_metrics.founders:
                        with tab_team:
                            for founder in metrics.team_metrics.founders:
                                st.markdown(f"**{founder.name}** - {founder.role}")
                                if founder.background: st.caption(f"📝 {founder.background}")
                                st.markdown("---")

                # REAL ESTATE PROFILE
                elif isinstance(metrics, REMetricsProfile):
                    tab_names = ["🏠 Finanziarie RE", "🌍 Mercato (VC)", "👥 Team"]
                    tab_re, tab_market, tab_team = st.tabs(tab_names)
                    with tab_re:
                        m = metrics.re_metrics
                        if m:
                            statuses = getattr(m, "metrics_status", {}) or {}
                            c1, c2 = st.columns(2)
                            with c1:
                                render_confidence_metric("Cap Rate", f"{m.cap_rate}%" if m.cap_rate else "-",
                                                         statuses.get("cap_rate", "UNVERIFIED"))
                                render_confidence_metric("Occupancy",
                                                         f"{m.occupancy_rate}%" if m.occupancy_rate else "-",
                                                         statuses.get("occupancy_rate", "UNVERIFIED"))
                            with c2:
                                render_confidence_metric("IRR", f"{m.irr}%" if m.irr else "-",
                                                         statuses.get("irr", "UNVERIFIED"))
                                render_confidence_metric("NOI",
                                                         f"${m.net_operating_income}M" if m.net_operating_income else "-",
                                                         statuses.get("net_operating_income", "UNVERIFIED"))
                        else:
                            st.info("Nessuna metrica RE.")

                # PHARMA PROFILE
                elif isinstance(metrics, PharmaMetricsProfile):
                    tab_names = ["🧪 R&D Pipeline", "👥 Team", "💵 Fundraising"]
                    tab_rd, tab_team, tab_fundraising = st.tabs(tab_names)
                    with tab_rd:
                        m = metrics.rd_metrics
                        if m:
                            statuses = getattr(m, "metrics_status", {}) or {}
                            c1, c2, c3 = st.columns(3)
                            with c1:
                                render_confidence_metric("Fase", m.clinical_trial_phase or "-",
                                                         statuses.get("clinical_trial_phase", "UNVERIFIED"))
                            with c2:
                                render_confidence_metric("FDA", m.fda_approval_status or "-",
                                                         statuses.get("fda_approval_status", "UNVERIFIED"))
                            with c3:
                                render_confidence_metric("Time to Market",
                                                         f"{m.time_to_market_years} y" if m.time_to_market_years else "-",
                                                         statuses.get("time_to_market_years", "UNVERIFIED"))
                            st.markdown("---")
                            c4, c5 = st.columns(2)
                            with c4:
                                render_confidence_metric("Brevetto", m.patent_expiry_date or "-",
                                                         statuses.get("patent_expiry_date", "UNVERIFIED"))
                            with c5:
                                render_confidence_metric("R&D Burn",
                                                         f"${m.rd_burn_rate_m}M" if m.rd_burn_rate_m else "-",
                                                         statuses.get("rd_burn_rate_m", "UNVERIFIED"))
                        else:
                            st.info("Nessuna metrica Pharma.")

                # LEGAL PROFILE
                elif isinstance(metrics, LegalMetricsProfile):
                    tab_names = ["⚖️ Risk & Compliance", "👥 Team"]
                    tab_legal, tab_team = st.tabs(tab_names)
                    with tab_legal:
                        m = metrics.legal_metrics
                        if m:
                            statuses = getattr(m, "metrics_status", {}) or {}
                            c1, c2 = st.columns(2)
                            gdpr_val = "✅ OK" if m.gdpr_compliance else "❌ NO"
                            with c1:
                                render_confidence_metric("GDPR", gdpr_val,
                                                         statuses.get("gdpr_compliance", "UNVERIFIED"))
                                render_confidence_metric("Contenziosi", str(m.pending_litigation_count),
                                                         statuses.get("pending_litigation_count", "UNVERIFIED"))
                            with c2:
                                render_confidence_metric("ISO/SOC", str(m.iso_soc_certified),
                                                         statuses.get("iso_soc_certified", "UNVERIFIED"))
                        else:
                            st.info("Nessuna metrica Legal.")

        # --- Report Testuali ---
        with results_placeholder.container():
            tab_summary, tab_risk, tab_feasibility, tab_competitors, tab_facts = st.tabs([
                "📋 Executive Summary", "🚩 Risk Analysis", "✅ Feasibility", "⚔️ Competitors", "🔍 Fact-Check"
            ])
            with tab_summary: st.markdown(result.get("executive_summary", ""))
            with tab_risk: st.markdown(result.get("risk_analysis", ""))
            with tab_feasibility: st.markdown(result.get("feasibility_analysis", ""))
            with tab_facts: render_fact_checking_table(result.get("fact_checking_table", []))
            with tab_competitors:
                comp_data = result.get("competitor_analysis")
                # Gestione robusta: riconverti da dict a oggetto se necessario (es. dopo ricaricamento sessione)
                if isinstance(comp_data, dict):
                    try:
                        comp_data = CompetitorAnalysis(**comp_data)
                    except Exception:
                        pass

                render_competitor_tab(comp_data)

            st.markdown("---")
            st.markdown(result.get("metrics_analysis", ""))

        # --- GRAFO (Ora persistente!) ---
        with graph_placeholder.container():
            if "graph_data" in result:
                render_knowledge_graph(
                    result["graph_data"]["nodes"],
                    result["graph_data"]["edges"],
                    entity_to_analyze
                )

    # --- SIDEBAR & DOWNLOADS ---
    if analysis_mode == "Matching Requisiti Tecnici":
        render_requirements_sidebar(result.get("executive_summary", ""), result.get("metadata", {}))
    else:
        render_metadata_sidebar(result.get("metadata", {}))

    if "source_documents" in result:
        render_source_inspector(result["source_documents"])

    st.markdown("---")
    st.markdown("### 📥 Download Report")

    comp_analysis = result.get('competitor_analysis')
    comp_text_section = "N/A"

    if comp_analysis:
        # Gestione dict/object
        if isinstance(comp_analysis, dict):
            c_obj = comp_analysis
        else:
            c_obj = comp_analysis.model_dump()

        if c_obj:
            comp_text_section = f"POSITIONING: {c_obj.get('market_position', 'N/A')}\n\n"
            for c in c_obj.get('competitors', []):
                # FIX: Aggiungiamo l'URL nel testo del report
                url_str = f" [Link: {c.get('website')}]" if c.get('website') else ""

                comp_text_section += f"- {c['name'].upper()}{url_str}\n"
                comp_text_section += f"  Desc: {c['description']}\n"
                comp_text_section += f"  Diff: {c.get('differentiation', '')}\n\n"

    # Aggiorna la stringa del report completo
    full_report_text = f"""REPORT DI ANALISI: {entity_to_analyze}
        --------------------------------------------------
        Settore: {sector}
        Data: {time.strftime("%Y-%m-%d %H:%M:%S")}

        1. EXECUTIVE SUMMARY
        --------------------
        {result.get('executive_summary', 'N/A')}

        2. COMPETITIVE LANDSCAPE
        ------------------------
        {comp_text_section}

        3. RISK ANALYSIS
        ----------------
        {result.get('risk_analysis', 'N/A')}

        4. FEASIBILITY ANALYSIS
        -----------------------
        {result.get('feasibility_analysis', 'N/A')}

        5. METRICS ANALYSIS
        -------------------
        {result.get('metrics_analysis', 'N/A')}
        """

    # Preparazione DataFrame Fact-Checking per CSV/Excel
    df_facts = pd.DataFrame(result.get("fact_checking_table", []))

    # Colonne per il download (4 pulsanti in fila)
    col_d1, col_d2, col_d3, col_d4 = st.columns(4)

    # --- A. DOWNLOAD TXT ---
    with col_d1:
        st.download_button(
            "📄 Report TXT",
            full_report_text,
            file_name=f"Report_{entity_to_analyze}.txt",
            mime="text/plain"
        )

    # --- B. DOWNLOAD CSV (Solo Fact Checking) ---
    with col_d2:
        if not df_facts.empty:
            csv_data = df_facts.to_csv(index=False).encode('utf-8')
            st.download_button(
                "📊 Fact-Check CSV",
                csv_data,
                file_name=f"FactCheck_{entity_to_analyze}.csv",
                mime="text/csv"
            )
        else:
            st.button("📊 CSV (No Dati)", disabled=True)

    # --- C. DOWNLOAD EXCEL (Report Completo) ---
    with col_d3:
        buffer_excel = io.BytesIO()
        with pd.ExcelWriter(buffer_excel, engine='xlsxwriter') as writer:
            # Foglio 1: Fact Checking
            if not df_facts.empty:
                df_facts.to_excel(writer, sheet_name='Fact Checking', index=False)

            # Foglio 2: Summary Testuale (Hack per mettere testo in celle)
            df_summary = pd.DataFrame([x.split('\n') for x in full_report_text.split('\n\n')])
            df_summary.to_excel(writer, sheet_name='Report Text', index=False, header=False)

        st.download_button(
            "📗 Report Excel",
            buffer_excel.getvalue(),
            file_name=f"Report_{entity_to_analyze}.xlsx",
            mime="application/vnd.ms-excel"
        )

    # --- D. DOWNLOAD PDF (Con sezione FONTI dedicata) ---
    with col_d4:
        class PDF(FPDF):
            def header(self):
                self.set_font('Arial', 'B', 16)
                self.cell(0, 10, f'VC Report: {entity_to_analyze}', 0, 1, 'C')
                self.set_font('Arial', 'I', 10)
                self.cell(0, 10, f'Generated by Spectre AI - {time.strftime("%Y-%m-%d")}', 0, 1, 'C')
                self.ln(5)

            def footer(self):
                self.set_y(-15)
                self.set_font('Arial', 'I', 8)
                self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

            def chapter_title(self, label):
                self.set_font('Arial', 'B', 14)
                self.set_fill_color(240, 240, 240)
                self.cell(0, 10, label, 0, 1, 'L', 1)
                self.ln(2)

            def replace_emojis(self, text):
                replacements = {
                    "✅": "[OK]", "❌": "[NO]", "⚠️": "[!]", "❓": "[?]",
                    "🚩": "[RISK]", "📋": "[SUMMARY]", "⚔️": "[VS]",
                    "🏢": "[AZIENDA]", "📊": "[METRICS]", "🔍": "[CHECK]",
                    "💰": "[FINANCE]", "🔧": "[TECH]", "🎯": "[TARGET]",
                    "👉": "->", "⭐": "*"
                }
                for emoji_char, replacement in replacements.items():
                    text = text.replace(emoji_char, replacement)
                return text

            def chapter_body(self, body):
                self.set_font('Arial', '', 11)
                clean_body = self.replace_emojis(body)
                clean_body = clean_body.replace('**', '').replace('## ', '').replace('### ', '')
                safe_body = clean_body.encode('latin-1', 'replace').decode('latin-1')
                self.multi_cell(0, 6, safe_body)
                self.ln(5)

            def add_link_line(self, label, url):
                """Aggiunge una riga con un link cliccabile."""
                self.set_font('Arial', '', 10)
                # Testo etichetta
                safe_label = label.encode('latin-1', 'replace').decode('latin-1')
                self.write(6, f"- {safe_label}: ")

                # Link blu e sottolineato
                self.set_text_color(0, 0, 255)
                self.set_font('Arial', 'U', 10)
                safe_url = url.encode('latin-1', 'replace').decode('latin-1')
                self.write(6, safe_url, link=url)

                # Reset stile
                self.set_text_color(0, 0, 0)
                self.set_font('Arial', '', 10)
                self.ln(6)


        try:
            pdf = PDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()

            # --- CAPITOLI STANDARD (1-5) ---
            pdf.chapter_title("1. EXECUTIVE SUMMARY")
            pdf.chapter_body(result.get('executive_summary', 'Nessun dato.'))

            pdf.chapter_title("2. COMPETITIVE LANDSCAPE")
            # Logica testo competitor (semplificata per il body)
            comp_analysis = result.get('competitor_analysis')
            comp_text_pdf = "Nessun dato.\n"
            if comp_analysis:
                if isinstance(comp_analysis, dict):
                    c_obj = comp_analysis
                else:
                    c_obj = comp_analysis.model_dump()
                if c_obj:
                    comp_text_pdf = f"POSIZIONAMENTO: {c_obj.get('market_position', 'N/A')}\n\n"
                    for i, c in enumerate(c_obj.get('competitors', []), 1):
                        comp_text_pdf += f"{i}. {c['name'].upper()}\n   {c['description']}\n\n"
            pdf.chapter_body(comp_text_pdf)

            pdf.chapter_title("3. RISK ANALYSIS")
            pdf.chapter_body(result.get('risk_analysis', 'Nessun dato.'))

            pdf.chapter_title("4. FEASIBILITY ANALYSIS")
            pdf.chapter_body(result.get('feasibility_analysis', 'Nessun dato.'))

            pdf.chapter_title("5. METRICS ANALYSIS")
            pdf.chapter_body(result.get('metrics_analysis', 'Nessun dato.'))

            # --- CAPITOLO 6: FONTI & RIFERIMENTI (NUOVO) ---
            pdf.add_page()  # Nuova pagina per le fonti
            pdf.chapter_title("6. SOURCES & REFERENCES")

            # 6.1 Link Competitor
            pdf.set_font('Arial', 'B', 11)
            pdf.cell(0, 8, "Competitor Websites:", 0, 1)

            has_sources = False
            if comp_analysis:
                if isinstance(comp_analysis, dict):
                    c_obj = comp_analysis
                else:
                    c_obj = comp_analysis.model_dump()
                for c in c_obj.get('competitors', []):
                    if c.get('website'):
                        pdf.add_link_line(c['name'], c['website'])
                        has_sources = True

            pdf.ln(4)

            # 6.2 Link Fact-Checking
            pdf.set_font('Arial', 'B', 11)
            pdf.cell(0, 8, "Fact-Checking Sources:", 0, 1)

            claims = result.get("fact_checking_table", [])
            for claim in claims:
                url = claim.get("source_url")
                if url:
                    # Usiamo l'inizio dell'affermazione come etichetta
                    label = claim.get("soggetto", "Fonte") + " (" + claim.get("status", "") + ")"
                    pdf.add_link_line(label, url)
                    has_sources = True

            if not has_sources:
                pdf.chapter_body("Nessun link esterno rilevato.")

            # Output
            pdf_output = pdf.output(dest='S').encode('latin-1')

            st.download_button(
                "📕 Report PDF (Full)",
                pdf_output,
                file_name=f"Report_{entity_to_analyze}.pdf",
                mime="application/pdf"
            )
        except Exception as e:
            st.error(f"Errore generazione PDF: {e}")

# Footer
st.markdown("---")
st.caption("🤖 Powered by Agentic KRAG | Built with LangChain, Neo4j, Streamlit")